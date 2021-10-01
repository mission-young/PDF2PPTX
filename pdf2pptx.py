# python 
# windows
# pyinstaller -i pdf2pptx.ico --add-data=".\default.pptx;." -F -w pdf2pptx.py 
# macOS
# pyinstaller -i pdf2pptx.ico --add-data=".\default.pptx:." -F  pdf2pptx.py 
# pipenv 下打包缩小体积

# pip install pyinstaller pymupdf python-pptx PyPDF2
import fitz # pip install  pymupdf
import os
import sys
import shutil
import pptx #pip install python-pptx
import PyPDF2


#生成资源文件目录访问路径
def resource_path(relative_path):
    if getattr(sys, 'frozen', False): #是否Bundle Resource
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

def pdf2pptx(pdf_name,zoom):
    # pdf2pngs
    pdf_name=os.path.abspath(pdf_name)
    if pdf_name[-4:].lower() == '.pdf':
        dir_name = os.path.dirname(pdf_name) # 获得地址的父链接

        base_name = os.path.basename(pdf_name)[0:-4] # 获得地址的文件名

        path = dir_name + os.sep + base_name

        if not os.path.exists(path):
            os.makedirs(path)

        with open(pdf_name,'rb') as f:
            pdfwidth=PyPDF2.PdfFileReader(f).getPage(0).mediaBox[2]
            pdfheight=PyPDF2.PdfFileReader(f).getPage(0).mediaBox[3]

        pdfscale=pdfheight/pdfwidth
        pdfmode = 0
        if pdfscale > 0.74 and pdfscale < 0.76:
            pdfmode = 0
        elif pdfscale > 0.5 and pdfscale < 0.6:
            pdfmode = 1
        else:
            pdfmode = -1

        pdf = fitz.open(pdf_name)
        # pdfwidth=pdf[0].get_images()[0][2]
        # pdfheight=pdf[0].get_images()[0][3]
        # pdfscale=pdfheight/pdfwidth
        # pdfmode = 0
        # if pdfscale > 0.74 and pdfscale < 0.76:
        #     pdfmode = 0
        # elif pdfscale > 0.5 and pdfscale < 0.6:
        #     pdfmode = 1
        # else:
        #     pdfmode = -1

        for pg in range(0, pdf.pageCount):
            page = pdf[pg]  # 获得每一页的对象
            trans = fitz.Matrix(zoom, zoom).preRotate(0)
            pm = page.getPixmap(matrix=trans, alpha=False)  # 获得每一页的流对象
            pm.writePNG(path + os.sep + base_name + '_' + '{:0>3d}.{}'.format(pg+1, 'png'))  # 保存图片
        pdf.close()

        print('PDF转PNG成功!')

    # pngs2pptx
        pngs=os.listdir(path+os.sep)
        template=resource_path(os.path.join("default.pptx"))
        prs=pptx.Presentation(template)
        if pdfmode == 0:
            prs.slide_width = pptx.util.Inches(4)
            prs.slide_height = pptx.util.Inches(3)
        elif pdfmode == 1:
            prs.slide_width = pptx.util.Inches(16)
            prs.slide_height = pptx.util.Inches(9)
        else:
            print('请采用标准格式!')
            prs.slide_width = pptx.util.Inches(4)
            prs.slide_height = pptx.util.Inches(3) 

        layout=prs.slide_layouts[6]
        for png in pngs:
            slide=prs.slides.add_slide(layout)
            slide.shapes.add_picture(path + os.sep + png,0,0,height=prs.slide_height)

        prs.save(path + '.pptx')
        print('PNG转PPTX成功!')
        shutil.rmtree(path)
    else:
        print('警告! 文件类型错误,请打开pdf文件类型!')

if __name__=="__main__":
    if len(sys.argv) == 2:
        pdf2pptx(sys.argv[1],5)
    elif len(sys.argv) > 2:
        pdf2pptx(sys.argv[1],int(sys.argv[2]))
    else:
        print('请输入要转化的文件名')